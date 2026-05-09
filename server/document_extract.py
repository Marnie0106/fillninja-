"""Extract plain text from uploaded PDF, Word (.docx), PowerPoint (.pptx)."""

from io import BytesIO
from pathlib import Path

MAX_UPLOAD_BYTES = 15 * 1024 * 1024
MAX_EXTRACT_CHARS = 400_000


def extract_text_from_bytes(data: bytes, filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if len(data) > MAX_UPLOAD_BYTES:
        raise ValueError(f"File too large ({len(data)} bytes). Maximum is {MAX_UPLOAD_BYTES // (1024 * 1024)} MB.")
    if suffix == ".pdf":
        return _extract_pdf(data)
    if suffix == ".docx":
        return _extract_docx(data)
    if suffix == ".pptx":
        return _extract_pptx(data)
    if suffix in (".ppt", ".doc"):
        raise ValueError(f"Legacy {suffix} is not supported. Save as {'docx' if suffix == '.doc' else 'pptx'} and upload again.")
    # Treat unknown as plain text
    try:
        return _trim_text(data.decode("utf-8", errors="replace"))
    except Exception:
        raise ValueError("Unsupported file type. Upload .pdf, .docx, .pptx.")


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        t = page.extract_text() or ""
        if t.strip():
            parts.append(t)
    return _trim_text("\n\n".join(parts).strip())


def _extract_docx(data: bytes) -> str:
    from docx import Document
    document = Document(BytesIO(data))
    parts: list[str] = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return _trim_text("\n".join(parts).strip())


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                t = (para.text or "").strip()
                if t:
                    parts.append(t)
    return _trim_text("\n".join(parts).strip())


def _trim_text(text: str) -> str:
    if len(text) <= MAX_EXTRACT_CHARS:
        return text
    return text[:MAX_EXTRACT_CHARS] + "\n\n[... document truncated for processing ...]"
